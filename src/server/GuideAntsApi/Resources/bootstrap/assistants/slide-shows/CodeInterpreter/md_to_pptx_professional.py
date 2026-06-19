
"""
Professional Markdown to PowerPoint Converter
Creates beautifully formatted PPTX presentations from Markdown with proper layouts, typography, and styling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import re
import os
import json
import tempfile
from PIL import Image
from urllib.request import urlopen, Request
from io import BytesIO


def _fix_presentation_xml(prs):
    """
    Fix common issues in presentation XML that cause PowerPoint to show a "repair" dialog.

    This function should be called AFTER all slides are fully created,
    right before saving the presentation.

    Issues fixed:
    1. Empty grpSpPr elements in slides - adds required xfrm transform element
    2. Mismatched slide size type attribute - removes type when custom dimensions are used

    Args:
        prs: A python-pptx Presentation object
    """
    # Fix 1: Add xfrm to grpSpPr in all slides
    for slide in prs.slides:
        # Access the slide's shape tree XML
        spTree = slide.shapes._spTree

        # Find the grpSpPr element
        grpSpPr = spTree.find(qn('p:grpSpPr'))
        if grpSpPr is not None:
            # Check if xfrm already exists
            xfrm = grpSpPr.find(qn('a:xfrm'))
            if xfrm is None:
                # Create the xfrm element with required child elements
                xfrm = etree.SubElement(grpSpPr, qn('a:xfrm'))

                # Add required child elements with zero values
                off = etree.SubElement(xfrm, qn('a:off'))
                off.set('x', '0')
                off.set('y', '0')

                ext = etree.SubElement(xfrm, qn('a:ext'))
                ext.set('cx', '0')
                ext.set('cy', '0')

                chOff = etree.SubElement(xfrm, qn('a:chOff'))
                chOff.set('x', '0')
                chOff.set('y', '0')

                chExt = etree.SubElement(xfrm, qn('a:chExt'))
                chExt.set('cx', '0')
                chExt.set('cy', '0')

    # Fix 2: Remove 'type' attribute from sldSz if present
    # python-pptx's default template has type="screen4x3" but when custom
    # dimensions are set (e.g., 16:9), the type attribute becomes invalid
    # and causes PowerPoint to show a repair dialog
    presentation_part = prs.part
    presentation_elm = presentation_part._element
    sldSz = presentation_elm.find(qn('p:sldSz'))
    if sldSz is not None and 'type' in sldSz.attrib:
        del sldSz.attrib['type']


# Default headers for URL requests - many servers block requests without User-Agent
DEFAULT_URL_HEADERS = {
    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
}


def _create_url_request(url):
    """Create a Request object with proper headers for downloading images"""
    return Request(url, headers=DEFAULT_URL_HEADERS)


# Supported image formats for python-pptx
PPTX_SUPPORTED_FORMATS = {'BMP', 'GIF', 'JPEG', 'JPG', 'PNG', 'TIFF', 'WMF'}


def _convert_to_supported_format(image_path):
    """
    Convert image to a format supported by python-pptx if needed.
    Returns the path to use (original or converted temp file), and whether cleanup is needed.
    """
    try:
        with Image.open(image_path) as img:
            img_format = img.format
            
            if img_format and img_format.upper() in PPTX_SUPPORTED_FORMATS:
                return image_path, False  # No conversion needed
            
            # Convert to PNG (lossless, good for most cases)
            
            # Handle different modes
            if img.mode in ('RGBA', 'LA', 'P'):
                # Keep transparency by converting to RGBA then PNG
                if img.mode == 'P':
                    img = img.convert('RGBA')
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Save to temp file
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            img.save(temp_file.name, 'PNG')
            temp_file.close()
            return temp_file.name, True  # Needs cleanup
            
    except Exception as e:
        return image_path, False  # Return original, let downstream handle errors


class ThemeLoader:
    """Loads theme configuration from JSON file"""
    
    @staticmethod
    def find_theme_file():
        """Automatically find theme file in common locations"""
        # Try script directory first (using __file__ which should be set by importer)
        try:
            # Access __file__ from global scope if available
            if '__file__' in globals():
                script_dir = os.path.dirname(os.path.realpath(globals()['__file__']))
                theme_path = os.path.join(script_dir, 'pptx_theme.json')
                if os.path.exists(theme_path):
                    return theme_path
        except:
            pass
            
        # Try current working directory
        cwd_theme = os.path.join(os.getcwd(), 'pptx_theme.json')
        if os.path.exists(cwd_theme):
            return cwd_theme
        
        # Try same directory as markdown file (if we know it)
        return None
    
    @staticmethod
    def load_theme(theme_file=None):
        """Load theme from JSON file, or use defaults if file not found"""
        # If no theme file specified, try to find it automatically
        if theme_file is None:
            theme_file = ThemeLoader.find_theme_file()
        
        default_theme = {
            "typography": {
                "fonts": {
                    "title": "Calibri Light",
                    "body": "Calibri",
                    "code": "Consolas"
                },
                "sizes": {
                    "title": 54,
                    "subtitle": 32,
                    "body": 20,
                    "bullet": 20,
                    "code": 16,
                    "caption": 16
                },
                "spacing": {
                    "title_space_after": 12,
                    "body_space_after": 6,
                    "bullet_space_after": 8,
                    "line_spacing": 1.2
                }
            },
            "colors": {
                "primary": [31, 56, 100],
                "secondary": [89, 89, 89],
                "accent": [0, 120, 212],
                "background": [255, 255, 255],
                "light_bg": [248, 249, 250],
                "code_bg": [248, 248, 248],
                "code_border": [200, 200, 200]
            },
            "layout": {
                "slide_dimensions": {
                    "width": 13.33,
                    "height": 7.5
                },
                "margins": {
                    "top": 0.5,
                    "bottom": 0.5,
                    "left": 1.0,
                    "right": 1.0
                },
                "spacing": {
                    "title_to_content": 0.75,
                    "paragraph_spacing": 0.25,
                    "bullet_spacing": 0.15,
                    "image_to_text": 0.5,
                    "bullet_indent_base": 0.5
                },
                "image": {
                    "max_width": 9.0,
                    "max_height": 4.5,
                    "default_max_height": 4.5
                }
            },
            "content": {
                "footer": None
            }
        }
        
        if theme_file and os.path.exists(theme_file):
            try:
                with open(theme_file, 'r', encoding='utf-8') as f:
                    theme_data = json.load(f)
                    
                    # Deep merge function for nested dictionaries
                    def deep_merge(default, override):
                        """Recursively merge override into default"""
                        result = default.copy()
                        for key, value in override.items():
                            if key in result and isinstance(result[key], dict) and isinstance(value, dict):
                                result[key] = deep_merge(result[key], value)
                            else:
                                result[key] = value
                        return result
                    
                    # Merge with defaults (theme data takes precedence)
                    if 'typography' in theme_data:
                        default_theme['typography'] = deep_merge(default_theme['typography'], theme_data['typography'])
                    if 'colors' in theme_data:
                        default_theme['colors'].update(theme_data['colors'])  # Colors are flat, simple update is fine
                    if 'layout' in theme_data:
                        default_theme['layout'] = deep_merge(default_theme['layout'], theme_data['layout'])
                    if 'content' in theme_data:
                        if default_theme.get('content') is None:
                            default_theme['content'] = theme_data['content']
                        else:
                            default_theme['content'] = deep_merge(default_theme['content'], theme_data['content'])
                    
                    # print(f"Theme loaded from: {theme_file}")
                    # print(f"  Title font: {default_theme['typography']['fonts']['title']}")
                    # print(f"  Primary color: {default_theme['colors']['primary']}")
                    return default_theme
            except Exception as e:
                print(f"Warning: Could not load theme file {theme_file}: {e}")
                print("Using default theme.")
                import traceback
                traceback.print_exc()
        else:
            if theme_file:
                print(f"Theme file not found: {theme_file}")
            # Don't print message if we auto-searched and didn't find it (that's OK)
        
        return default_theme


class Typography:
    """Typography presets for consistent styling - now theme-aware"""
    
    _theme = None
    
    @classmethod
    def load_theme(cls, theme_file=None):
        """Load theme configuration"""
        cls._theme = ThemeLoader.load_theme(theme_file)
    
    @classmethod
    def _get_font(cls, font_type):
        """Get font name from theme"""
        if cls._theme:
            return cls._theme['typography']['fonts'].get(font_type, 'Calibri')
        # Fallback defaults
        defaults = {'title': 'Calibri Light', 'body': 'Calibri', 'code': 'Consolas'}
        return defaults.get(font_type, 'Calibri')
    
    @classmethod
    def _get_size(cls, size_type):
        """Get font size from theme"""
        if cls._theme:
            return Pt(cls._theme['typography']['sizes'].get(size_type, 20))
        # Fallback defaults
        defaults = {'title': 54, 'subtitle': 32, 'body': 20, 'bullet': 20, 'code': 16, 'caption': 16}
        return Pt(defaults.get(size_type, 20))
    
    @classmethod
    def _get_color(cls, color_type):
        """Get color from theme"""
        if cls._theme:
            rgb = cls._theme['colors'].get(color_type, [31, 56, 100])
            return RGBColor(rgb[0], rgb[1], rgb[2])
        # Fallback defaults
        defaults = {
            'primary': [31, 56, 100],
            'secondary': [89, 89, 89],
            'accent': [0, 120, 212],
            'background': [255, 255, 255],
            'light_bg': [248, 249, 250],
            'code_bg': [248, 248, 248],
            'code_border': [200, 200, 200]
        }
        rgb = defaults.get(color_type, [31, 56, 100])
        return RGBColor(rgb[0], rgb[1], rgb[2])
    
    @classmethod
    def _get_spacing(cls, spacing_type):
        """Get spacing value from theme"""
        if cls._theme:
            return cls._theme['typography']['spacing'].get(spacing_type, 1.2)
        # Fallback defaults
        defaults = {'title_space_after': 12, 'body_space_after': 6, 'bullet_space_after': 8, 'line_spacing': 1.2}
        return defaults.get(spacing_type, 1.2)
    
    @classmethod
    def apply_title_style(cls, paragraph):
        """Apply title styling to a paragraph"""
        paragraph.font.name = cls._get_font('title')
        paragraph.font.size = cls._get_size('title')
        paragraph.font.bold = True
        paragraph.font.color.rgb = cls._get_color('primary')
        paragraph.space_after = Pt(cls._get_spacing('title_space_after'))
    
    @classmethod
    def apply_body_style(cls, paragraph, size=None):
        """Apply body text styling"""
        paragraph.font.name = cls._get_font('body')
        paragraph.font.size = size or cls._get_size('body')
        paragraph.font.color.rgb = cls._get_color('primary')
        paragraph.line_spacing = cls._get_spacing('line_spacing')
        paragraph.space_after = Pt(cls._get_spacing('body_space_after'))
    
    @classmethod
    def apply_code_style(cls, paragraph):
        """Apply code/monospace styling"""
        paragraph.font.name = cls._get_font('code')
        paragraph.font.size = cls._get_size('code')
        paragraph.font.color.rgb = cls._get_color('primary')
    
    # Class properties for backward compatibility (accessed as Typography.FONT_TITLE)
    @classmethod
    def get_font_title(cls):
        return cls._get_font('title')
    
    @classmethod
    def get_font_body(cls):
        return cls._get_font('body')
    
    @classmethod
    def get_font_code(cls):
        return cls._get_font('code')
    
    @classmethod
    def get_size_title(cls):
        return cls._get_size('title')
    
    @classmethod
    def get_size_subtitle(cls):
        return cls._get_size('subtitle')
    
    @classmethod
    def get_size_body(cls):
        return cls._get_size('body')
    
    @classmethod
    def get_size_bullet(cls):
        return cls._get_size('bullet')
    
    @classmethod
    def get_size_code(cls):
        return cls._get_size('code')
    
    @classmethod
    def get_size_caption(cls):
        return cls._get_size('caption')
    
    @classmethod
    def get_color_primary(cls):
        return cls._get_color('primary')
    
    @classmethod
    def get_color_secondary(cls):
        return cls._get_color('secondary')
    
    @classmethod
    def get_color_accent(cls):
        return cls._get_color('accent')
    
    @classmethod
    def get_color_background(cls):
        return cls._get_color('background')
    
    @classmethod
    def get_color_light_bg(cls):
        return cls._get_color('light_bg')
    
    @classmethod
    def get_color_code_bg(cls):
        return cls._get_color('code_bg')
    
    @classmethod
    def get_color_code_border(cls):
        return cls._get_color('code_border')


class LayoutEngine:
    """Handles slide layout and positioning - now theme-aware"""
    
    _theme = None
    
    @classmethod
    def load_theme(cls, theme_file=None):
        """Load theme configuration"""
        cls._theme = ThemeLoader.load_theme(theme_file)
    
    @classmethod
    def _get_slide_width(cls):
        """Get slide width from theme"""
        if cls._theme:
            return Inches(cls._theme['layout']['slide_dimensions'].get('width', 13.33))
        return Inches(13.33)
    
    @classmethod
    def _get_slide_height(cls):
        """Get slide height from theme"""
        if cls._theme:
            return Inches(cls._theme['layout']['slide_dimensions'].get('height', 7.5))
        return Inches(7.5)
    
    @classmethod
    def _get_margin(cls, margin_type):
        """Get margin from theme"""
        if cls._theme:
            return Inches(cls._theme['layout']['margins'].get(margin_type, 0.5))
        defaults = {'top': 0.5, 'bottom': 0.5, 'left': 1.0, 'right': 1.0}
        return Inches(defaults.get(margin_type, 0.5))
    
    @classmethod
    def _get_spacing(cls, spacing_type):
        """Get spacing from theme"""
        if cls._theme:
            return Inches(cls._theme['layout']['spacing'].get(spacing_type, 0.5))
        defaults = {
            'title_to_content': 0.75,
            'paragraph_spacing': 0.25,
            'bullet_spacing': 0.15,
            'image_to_text': 0.5,
            'bullet_indent_base': 0.5
        }
        return Inches(defaults.get(spacing_type, 0.5))
    
    @classmethod
    def _get_image_max(cls, dimension_type):
        """Get image max dimension from theme"""
        if cls._theme:
            return Inches(cls._theme['layout']['image'].get(dimension_type, 4.5))
        defaults = {'max_width': 9.0, 'max_height': 4.5, 'default_max_height': 4.5}
        return Inches(defaults.get(dimension_type, 4.5))
    
    # Class properties for backward compatibility
    @classmethod
    def get_slide_width(cls):
        return cls._get_slide_width()
    
    @classmethod
    def get_slide_height(cls):
        return cls._get_slide_height()
    
    @classmethod
    def get_margin_top(cls):
        return cls._get_margin('top')
    
    @classmethod
    def get_margin_bottom(cls):
        return cls._get_margin('bottom')
    
    @classmethod
    def get_margin_left(cls):
        return cls._get_margin('left')
    
    @classmethod
    def get_margin_right(cls):
        return cls._get_margin('right')
    
    @classmethod
    def get_footer_config(cls):
        """Get footer configuration from theme"""
        if cls._theme and 'content' in cls._theme:
            return cls._theme['content'].get('footer')
        return None
    
    def __init__(self, presentation):
        self.prs = presentation
        self.blank_layout = self.prs.slide_layouts[6]  # Blank layout
        
        # Calculate content area from theme
        self.SLIDE_WIDTH = self._get_slide_width()
        self.SLIDE_HEIGHT = self._get_slide_height()
        self.MARGIN_TOP = self._get_margin('top')
        self.MARGIN_BOTTOM = self._get_margin('bottom')
        self.MARGIN_LEFT = self._get_margin('left')
        self.MARGIN_RIGHT = self._get_margin('right')
        
        # Content area
        self.CONTENT_WIDTH = self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT
        self.CONTENT_HEIGHT = self.SLIDE_HEIGHT - self.MARGIN_TOP - self.MARGIN_BOTTOM
        self.CONTENT_LEFT = self.MARGIN_LEFT
        self.CONTENT_TOP = self.MARGIN_TOP
    
    def create_blank_slide(self):
        """Create a new blank slide"""
        return self.prs.slides.add_slide(self.blank_layout)
    
    def add_title(self, slide, text, top=None, center=True):
        """Add a title text box to slide"""
        if top is None:
            top = self.CONTENT_TOP
        
        left = self.CONTENT_LEFT
        width = self.CONTENT_WIDTH
        # Increase height to accommodate large title font (54pt) which may wrap
        height = Inches(1.5)
        
        if center:
            # Center the title
            text_box = slide.shapes.add_textbox(
                left, top, width, height
            )
            text_frame = text_box.text_frame
            # Use raw text for title for now (or process formatting too?)
            # Titles usually don't have bold markers but why not support it
            paragraph = text_frame.paragraphs[0]
            paragraph.text = "" # Clear default
            self._process_text_with_formatting(paragraph, text)
            
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            paragraph.alignment = PP_ALIGN.CENTER
            Typography.apply_title_style(paragraph)
        else:
            # Left-aligned title
            text_box = slide.shapes.add_textbox(
                left, top, width, height
            )
            text_frame = text_box.text_frame
            
            paragraph = text_frame.paragraphs[0]
            paragraph.text = "" # Clear default
            self._process_text_with_formatting(paragraph, text)
            
            text_frame.word_wrap = True
            
            paragraph.alignment = PP_ALIGN.LEFT
            Typography.apply_title_style(paragraph)
        
        return text_box
    
    def _process_text_with_formatting(self, paragraph, text):
        """Parse text for markdown formatting (bold, italic, bold+italic, links) and add runs"""
        # Process links first, then handle formatting in non-link segments
        link_pattern = r'\[([^\]]+)\]\(([^)]+)\)'
        
        last_end = 0
        for match in re.finditer(link_pattern, text):
            # Process text before this link (may contain formatting)
            before_text = text[last_end:match.start()]
            if before_text:
                self._add_formatted_runs(paragraph, before_text)
            
            # Add the link as a single run with hyperlink
            link_text = match.group(1)
            link_url = match.group(2)
            
            # Strip any formatting markers from link text for display
            display_text = link_text
            is_bold = False
            is_italic = False
            if display_text.startswith('***') and display_text.endswith('***') and len(display_text) >= 6:
                display_text = display_text[3:-3]
                is_bold = is_italic = True
            elif display_text.startswith('**') and display_text.endswith('**') and len(display_text) >= 4:
                display_text = display_text[2:-2]
                is_bold = True
            elif display_text.startswith('*') and display_text.endswith('*') and len(display_text) >= 2:
                display_text = display_text[1:-1]
                is_italic = True
            
            run = paragraph.add_run()
            run.text = display_text
            if is_bold:
                run.font.bold = True
            if is_italic:
                run.font.italic = True
            
            # Add hyperlink - use try/except to handle any issues gracefully
            try:
                run.hyperlink.address = link_url
            except Exception as e:
                print(f"Warning: Could not add hyperlink for '{link_url}': {e}")
            
            last_end = match.end()
        
        # Process any remaining text after the last link
        remaining_text = text[last_end:]
        if remaining_text:
            self._add_formatted_runs(paragraph, remaining_text)
    
    def _add_formatted_runs(self, paragraph, text):
        """Add runs with bold/italic formatting (no links)"""
        # Split by markdown formatting markers
        parts = re.split(r'(\*\*\*.*?\*\*\*|\*\*.*?\*\*|\*.*?\*)', text)
        
        for part in parts:
            if not part:
                continue
            
            run = paragraph.add_run()
            
            if part.startswith('***') and part.endswith('***') and len(part) >= 6:
                run.text = part[3:-3]
                run.font.bold = True
                run.font.italic = True
            elif part.startswith('**') and part.endswith('**') and len(part) >= 4:
                run.text = part[2:-2]
                run.font.bold = True
            elif part.startswith('*') and part.endswith('*') and len(part) >= 2:
                run.text = part[1:-1]
                run.font.italic = True
            else:
                run.text = part
    
    def add_text_box(self, slide, text, left, top, width, height, 
                     font_size=None, alignment=PP_ALIGN.LEFT):
        """Add a formatted text box"""
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        paragraph = text_frame.paragraphs[0]
        paragraph.text = "" # Clear default
        
        # Apply style first to set defaults for runs
        paragraph.alignment = alignment
        Typography.apply_body_style(paragraph, font_size)
        
        # Process formatting
        self._process_text_with_formatting(paragraph, text)
        
        return text_box
    
    def add_bullet_list(self, slide, items, left, top, width, height, 
                       default_indent_level=0):
        """Add a bullet list with proper formatting"""
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        # Clear default paragraph
        text_frame.clear()
        
        indent_base = self._get_spacing('bullet_indent_base')
        
        for i, item in enumerate(items):
            p = text_frame.add_paragraph()
            
            # Handle both string items and dict items with level
            if isinstance(item, dict) and 'text' in item:
                text = item['text']
                level = item.get('level', default_indent_level)
            else:
                text = str(item)
                level = default_indent_level
            
            # Apply style to paragraph first (sets defaults for runs)
            p.level = level
            p.font.name = Typography._get_font('body')
            p.font.size = Typography._get_size('bullet')
            p.font.color.rgb = Typography._get_color('primary')
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(Typography._get_spacing('bullet_space_after'))
            
            # Prepend bullet char if not relying on auto-bullets
            if not text.lstrip().startswith('•') and not text.lstrip().startswith('-'):
                # Add bullet as first run
                run = p.add_run()
                run.text = "• "
            
            # Process formatting in text
            self._process_text_with_formatting(p, text)
        
        return text_box
    
    def get_image_orientation(self, image_path):
        """Get image orientation: 'portrait', 'landscape', or 'square'"""
        # Handle URLs - download temporarily to determine orientation
        is_url = image_path.startswith('http://') or image_path.startswith('https://')
        temp_file = None
        
        if is_url:
            try:
                # Download image and check orientation directly from memory
                req = _create_url_request(image_path)
                with urlopen(req) as response:
                    img_data = response.read()
                    
                    # Check if we got HTML instead of image
                    if img_data[:50].lower().startswith(b'<!doctype') or img_data[:50].lower().startswith(b'<html'):
                        return 'landscape'
                    
                    # Check orientation directly from memory (no temp file needed)
                    from io import BytesIO
                    with Image.open(BytesIO(img_data)) as img:
                        img_width, img_height = img.size
                        aspect_ratio = img_width / img_height
                        
                        if 0.9 <= aspect_ratio <= 1.1:
                            return 'square'
                        elif aspect_ratio > 1.1:
                            return 'landscape'
                        else:
                            return 'portrait'
            except Exception as e:
                # Default to landscape for URLs if download fails
                return 'landscape'
        
        if not os.path.exists(image_path):
            if temp_file:
                try:
                    os.unlink(temp_file.name)
                except:
                    pass
            return None
        
        try:
            with Image.open(image_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height
                
                # Define thresholds for square (0.9 to 1.1)
                if 0.9 <= aspect_ratio <= 1.1:
                    result = 'square'
                elif aspect_ratio > 1.1:
                    result = 'landscape'
                else:
                    result = 'portrait'
                
                # Clean up temporary file
                if temp_file:
                    try:
                        os.unlink(temp_file.name)
                    except:
                        pass
                
                return result
        except Exception as e:
            print(f"Error reading image {image_path}: {e}")
            if temp_file:
                try:
                    os.unlink(temp_file.name)
                except:
                    pass
            return None
    
    def add_image(self, slide, image_path, left, top, max_width=None, 
                  max_height=None, maintain_aspect=True, width=None, height=None):
        """Add an image with smart sizing"""
        # Check if it's a URL
        is_url = image_path.startswith('http://') or image_path.startswith('https://')
        
        # For URLs, download to temporary file
        temp_file = None
        actual_image_path = image_path
        
        if is_url:
            try:
                # Download image to temporary file with proper headers
                url_request = _create_url_request(image_path)
                with urlopen(url_request) as req:
                    # Check HTTP status
                    if hasattr(req, 'status') and req.status != 200:
                        return None
                    
                    # Check content type to determine file extension
                    content_type = req.headers.get('Content-Type', 'image/jpeg')
                    
                    if 'png' in content_type.lower():
                        suffix = '.png'
                    elif 'gif' in content_type.lower():
                        suffix = '.gif'
                    elif 'webp' in content_type.lower():
                        suffix = '.webp'
                    elif 'jpeg' in content_type.lower() or 'jpg' in content_type.lower():
                        suffix = '.jpg'
                    else:
                        # Try to determine from URL extension
                        url_lower = image_path.lower()
                        if url_lower.endswith('.png'):
                            suffix = '.png'
                        elif url_lower.endswith('.gif'):
                            suffix = '.gif'
                        elif url_lower.endswith('.webp'):
                            suffix = '.webp'
                        else:
                            suffix = '.jpg'  # Default to jpg
                    
                    img_data = req.read()
                
                if not img_data:
                    return None
                
                # Check if we got HTML instead of image (common error response)
                if img_data[:100].lower().startswith(b'<!doctype') or img_data[:100].lower().startswith(b'<html'):
                    return None
                
                # Check if format needs conversion (WEBP not supported by python-pptx)
                try:
                    from io import BytesIO
                    with Image.open(BytesIO(img_data)) as img:
                        img_format = img.format
                        
                        if img_format and img_format.upper() not in PPTX_SUPPORTED_FORMATS:
                            # Convert to PNG
                            output = BytesIO()
                            # Handle transparency
                            if img.mode in ('RGBA', 'LA', 'P'):
                                if img.mode == 'P' and 'transparency' in img.info:
                                    img = img.convert('RGBA')
                                elif img.mode == 'P':
                                    img = img.convert('RGB')
                            elif img.mode not in ('RGB', 'RGBA'):
                                img = img.convert('RGB')
                            img.save(output, 'PNG')
                            img_data = output.getvalue()
                            suffix = '.png'
                except Exception as conv_err:
                    pass  # Continue with original format
                
                # Create temporary file with appropriate extension
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
                temp_file.write(img_data)
                temp_file.flush()  # Ensure data is written
                temp_file.close()
                actual_image_path = temp_file.name
                
                # Verify file was created and has content
                if not os.path.exists(actual_image_path):
                    print(f"Error: Temporary file was not created: {actual_image_path}")
                    return None
                
                file_size = os.path.getsize(actual_image_path)
                if file_size == 0:
                    print(f"Error: Temporary file is empty: {actual_image_path}")
                    try:
                        os.unlink(actual_image_path)
                    except:
                        pass
                    return None
                
                # Verify it's actually an image by trying to open it with PIL
                try:
                    with Image.open(actual_image_path) as test_img:
                        test_img.verify()  # Verify it's a valid image
                except Exception as img_error:
                    try:
                        os.unlink(actual_image_path)
                    except:
                        pass
                    return None
            except Exception as e:
                if temp_file and 'temp_file' in locals():
                    try:
                        if os.path.exists(temp_file.name):
                            os.unlink(temp_file.name)
                    except:
                        pass
                return None
        
        if not os.path.exists(actual_image_path):
            print(f"Warning: Image not found: {actual_image_path}")
            if temp_file:
                try:
                    os.unlink(temp_file.name)
                except:
                    pass
            return None
        
        # Get image dimensions
        aspect_ratio = None
        img_width_px = None
        try:
            with Image.open(actual_image_path) as img:
                img_width_px, img_height_px = img.size
                aspect_ratio = img_width_px / img_height_px
        except Exception as e:
            print(f"Error reading image {actual_image_path}: {e}")
            if temp_file:
                try:
                    os.unlink(temp_file.name)
                except:
                    pass
            return None
        
        # Use provided dimensions if specified, otherwise calculate
        if width is not None and height is not None:
            final_width = width
            final_height = height
        else:
            # Calculate dimensions
            if max_width is None:
                max_width = self.CONTENT_WIDTH
            if max_height is None:
                max_height = self._get_image_max('default_max_height')
            
            if maintain_aspect and aspect_ratio:
                # Calculate size maintaining aspect ratio
                if img_width_px:
                    calc_width = min(max_width, Inches(img_width_px / 96))  # 96 DPI assumption
                else:
                    calc_width = max_width
                calc_height = calc_width / aspect_ratio
                
                if calc_height > max_height:
                    calc_height = max_height
                    calc_width = calc_height * aspect_ratio
                
                final_width = calc_width
                final_height = calc_height
            else:
                final_width = max_width
                final_height = max_height
        
        try:
            # For local files (not URLs), check if conversion is needed
            if not is_url:
                converted_path, needs_conversion_cleanup = _convert_to_supported_format(actual_image_path)
                if needs_conversion_cleanup:
                    actual_image_path = converted_path
            else:
                needs_conversion_cleanup = False
            
            picture = slide.shapes.add_picture(actual_image_path, left, top, 
                                             width=final_width, height=final_height)
            
            # Cleanup converted temp file if we created one (for local files)
            if needs_conversion_cleanup:
                try:
                    os.unlink(actual_image_path)
                except:
                    pass
            
            # Clean up temporary file if we created one
            if temp_file:
                try:
                    os.unlink(temp_file.name)
                except:
                    pass  # Best effort cleanup
            
            return picture
        except Exception as e:
            error_msg = f"Error adding image {image_path}"
            if 'actual_image_path' in locals() and actual_image_path != image_path:
                error_msg += f" (temp file: {actual_image_path})"
            error_msg += f": {e}"
            print(error_msg)
            import traceback
            traceback.print_exc()
            # Clean up temporary file on error
            if 'temp_file' in locals() and temp_file:
                try:
                    if os.path.exists(temp_file.name):
                        os.unlink(temp_file.name)
                except:
                    pass
            return None
    
    def center_horizontally(self, width):
        """Calculate left position to center an element"""
        return (self.SLIDE_WIDTH - width) / 2
    
    def center_vertically(self, height):
        """Calculate top position to center an element"""
        return (self.SLIDE_HEIGHT - height) / 2
    
    def add_footer(self, slide, text):
        """Add footer text to slide"""
        if not text:
            return
            
        left = self.MARGIN_LEFT
        width = self.CONTENT_WIDTH
        height = Inches(0.3)
        top = self.SLIDE_HEIGHT - Inches(0.4) # Slightly above bottom edge
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = Typography.get_font_body()
        p.font.size = Pt(10)
        p.font.color.rgb = Typography.get_color_secondary()


class MarkdownParser:
    """Enhanced Markdown parser for slide content"""
    
    def __init__(self, md_text):
        self.md_text = md_text
        self.lines = md_text.splitlines()
        self.slides = []
    
    def parse(self):
        """Parse Markdown into slide structures"""
        current_slide = {
            'title': None,
            'subtitle': None,
            'content': [],
            'images': [],
            'type': None
        }
        
        in_code_block = False
        code_block_lines = []
        in_table = False
        table_rows = []
        table_headers = None
        
        for line in self.lines:
            line_stripped = line.strip()
            
            # Handle code blocks
            if line_stripped.startswith('```'):
                if in_code_block:
                    # End of code block
                    current_slide['content'].append({
                        'type': 'code',
                        'text': '\n'.join(code_block_lines)
                    })
                    code_block_lines = []
                    in_code_block = False
                else:
                    # Start of code block
                    in_code_block = True
                continue
            
            if in_code_block:
                code_block_lines.append(line)
                continue
            
            # Horizontal rule = slide break
            if line_stripped == '---' or line_stripped.startswith('---'):
                if current_slide['title'] is not None or current_slide['content'] or current_slide['images']:
                    self.slides.append(current_slide)
                current_slide = {
                    'title': None,
                    'subtitle': None,
                    'content': [],
                    'images': [],
                    'type': None
                }
                continue
            
            # Headers
            if line_stripped.startswith('# '):
                if current_slide['title'] is not None or current_slide['content'] or current_slide['images']:
                    self.slides.append(current_slide)
                current_slide = {
                    'title': line_stripped[2:].strip(),
                    'subtitle': None,
                    'content': [],
                    'images': [],
                    'type': None
                }
            elif line_stripped.startswith('## '):
                current_slide['subtitle'] = line_stripped[3:].strip()
            elif line_stripped.startswith('### '):
                current_slide['content'].append({
                    'type': 'heading3',
                    'text': line_stripped[4:].strip()
                })
            
            # Bullet lists
            elif line_stripped.startswith('- ') or line_stripped.startswith('* '):
                bullet_text = line_stripped[2:].strip()
                # Check for nested bullets
                indent_level = 0
                if line.startswith('  ') or line.startswith('	'):
                    indent_level = 1
                
                # Check if we need to start a new list
                # Start a new list if:
                # 1. There's no content yet, OR
                # 2. The last content item is NOT a bullets list (e.g., paragraph, heading, etc.)
                # This ensures that text between bullet lists creates separate lists
                if (not current_slide['content'] or 
                    current_slide['content'][-1]['type'] != 'bullets'):
                    current_slide['content'].append({
                        'type': 'bullets',
                        'items': []
                    })
                
                current_slide['content'][-1]['items'].append({
                    'text': bullet_text,
                    'level': indent_level
                })
            
            # Numbered lists
            elif re.match(r'^\d+\.\s', line_stripped):
                match = re.match(r'^\d+\.\s(.+)', line_stripped)
                if match:
                    item_text = match.group(1)
                    # Check if we need to start a new list
                    # Start a new list if:
                    # 1. There's no content yet, OR
                    # 2. The last content item is NOT a numbered list (e.g., paragraph, heading, etc.)
                    # This ensures that text between numbered lists creates separate lists
                    if (not current_slide['content'] or 
                        current_slide['content'][-1]['type'] != 'numbered'):
                        current_slide['content'].append({
                            'type': 'numbered',
                            'items': []
                        })
                    
                    current_slide['content'][-1]['items'].append(item_text)
            
            # Table rows - parse table structure
            elif line_stripped.startswith('|') and '|' in line_stripped[1:]:
                # Check if this is a separator row (|---|---|)
                is_separator = re.match(r'^\|[\s\-\|:]+\|$', line_stripped)
                if is_separator:
                    # This is a separator row - mark that we have headers if this is the first separator
                    if in_table and table_headers is None and table_rows:
                        table_headers = table_rows[0]
                        table_rows = table_rows[1:]  # Remove header from rows
                    continue
                
                # Parse table row
                cells = [cell.strip() for cell in line_stripped.split('|')[1:-1]]  # Remove empty first/last
                
                if not in_table:
                    # Starting a new table
                    in_table = True
                    table_rows = []
                    table_headers = None
                
                # Extract images from cells and clean cell text
                row_data = []
                for cell in cells:
                    cell_text = cell
                    cell_images = []
                    
                    # Extract images from cell
                    if '![' in cell:
                        # Use pattern that handles URLs with parentheses
                        image_pattern = r'!\[[^\]]*\]\(([^()\s]+(?:\([^)]*\)[^()\s]*)*)\)'
                        matches = re.findall(image_pattern, cell)
                        for image_path in matches:
                            # Clean up image path - preserve query params for URLs
                            original_path = image_path
                            if not image_path.startswith('http://') and not image_path.startswith('https://'):
                                image_path = image_path.split('?')[0]  # Remove query params for local files
                            if image_path.startswith('./Output/'):
                                image_path = image_path[len('./Output/'):]
                            cell_images.append(image_path)
                            current_slide['images'].append(image_path)
                        
                        # Remove image syntax from cell text (use same pattern)
                        cell_text = re.sub(r'!\[[^\]]*\]\([^()\s]+(?:\([^)]*\)[^()\s]*)*\)', '', cell).strip()
                    
                    row_data.append({
                        'text': cell_text,
                        'images': cell_images
                    })
                
                table_rows.append(row_data)
                continue
            
            # End of table (non-table line after table)
            elif in_table:
                # Save the table to current slide
                if table_rows:
                    current_slide['content'].append({
                        'type': 'table',
                        'headers': table_headers,
                        'rows': table_rows
                    })
                    # print(f"Added table with {len(table_rows)} rows")
                in_table = False
                table_rows = []
                table_headers = None
            
            # Images - search anywhere in the line (not just at start) to catch images
            elif '![' in line_stripped:
                # Find all image patterns in the line
                # Use a pattern that handles URLs with parentheses (common in CDN URLs)
                # This pattern allows for balanced parentheses within the URL
                image_pattern = r'!\[[^\]]*\]\(([^()\s]+(?:\([^)]*\)[^()\s]*)*)\)'
                matches = re.findall(image_pattern, line_stripped)
                for image_path in matches:
                    # Clean up image path - but preserve query params for URLs!
                    # Only strip query params for local files, not URLs
                    if not image_path.startswith('http://') and not image_path.startswith('https://'):
                        image_path = image_path.split('?')[0]  # Remove query params for local files
                    if image_path.startswith('./Output/'):
                        image_path = image_path[len('./Output/'):] 
                    current_slide['images'].append(image_path)
                
                # Remove image syntax from line to get remaining text (use same pattern that handles parens)
                line_without_images = re.sub(r'!\[[^\]]*\]\([^()\s]+(?:\([^)]*\)[^()\s]*)*\)', '', line_stripped).strip()
                # If there's remaining text after removing images, add it as content
                if line_without_images and not line_without_images.startswith('|'):
                    current_slide['content'].append({
                        'type': 'paragraph',
                        'text': line_without_images
                    })
                # Continue to next line (images are extracted, content handled if any)
                continue
            
            # Regular paragraphs
            elif line_stripped and not line_stripped.startswith('#'):
                current_slide['content'].append({
                    'type': 'paragraph',
                    'text': line_stripped
                })
        
        # End any open table
        if in_table and table_rows:
            current_slide['content'].append({
                'type': 'table',
                'headers': table_headers,
                'rows': table_rows
            })
            # print(f"Added table with {len(table_rows)} rows (end of file)")
        
        # Add last slide (include slides with just images, even if no other content)
        if current_slide['title'] is not None or current_slide['content'] or current_slide['images']:
            self.slides.append(current_slide)
        
        return self.slides
    
    def detect_slide_type(self, slide_data):
        """Detect the best layout type for a slide"""
        # Check for title slides (H1 only or H1+H2 only with one image)
        has_title = slide_data['title'] is not None
        has_subtitle = slide_data['subtitle'] is not None
        has_one_image = len(slide_data['images']) == 1
        has_no_content = not slide_data['content'] or all(
            c['type'] in ['heading3', 'heading4'] for c in slide_data['content']
        )
        
        # Title slide: H1 only with one image
        if has_title and not has_subtitle and has_one_image and has_no_content:
            return 'title_slide'
        
        # Title slide with subtitle: H1 and H2 only with one image
        if has_title and has_subtitle and has_one_image and has_no_content:
            return 'title_slide_subtitle'
        
        # Other slide types
        # Check if slide has images
        has_images = slide_data.get('images') and len(slide_data['images']) > 0
        has_content = slide_data.get('content') and len(slide_data['content']) > 0
        
        # Check for specific content types
        has_table = any(c['type'] == 'table' for c in slide_data['content'])
        has_code = any(c['type'] == 'code' for c in slide_data['content'])
        
        if has_table:
            return 'table'
        elif has_code:
            return 'code'
        elif has_images and not has_content:
            return 'image_focused'
        elif has_images and has_content:
            return 'image_content'
        elif any(c['type'] == 'bullets' for c in slide_data['content']) or any(c['type'] == 'numbered' for c in slide_data['content']):
            return 'bullets'
        elif slide_data['subtitle']:
            return 'title_subtitle'
        else:
            return 'title_content'


class PowerPointGenerator:
    """Main generator class"""
    
    def __init__(self, md_file, output_file, theme_file=None):
        self.md_file = md_file
        self.output_file = output_file
        
        # Store markdown directory for resolving relative image paths
        if os.path.exists(md_file):
            self.md_dir = os.path.dirname(os.path.abspath(md_file))
        else:
            self.md_dir = os.getcwd()
        
        # Auto-find theme file if not specified
        if theme_file is None:
            # Try same directory as markdown file first
            if os.path.exists(md_file):
                md_theme = os.path.join(self.md_dir, 'pptx_theme.json')
                if os.path.exists(md_theme):
                    theme_file = md_theme
            
            # Fall back to auto-search (script dir, cwd, etc.)
            if theme_file is None:
                theme_file = ThemeLoader.find_theme_file()
        
        self.theme_file = theme_file
        
        # Load theme before creating presentation
        Typography.load_theme(theme_file)
        LayoutEngine.load_theme(theme_file)
        
        self.prs = Presentation()
        self.layout_engine = LayoutEngine(self.prs)
        self.prs.slide_width = self.layout_engine.SLIDE_WIDTH
        self.prs.slide_height = self.layout_engine.SLIDE_HEIGHT
        
        self.parser = None
    
    def _resolve_image_path(self, image_path):
        """Resolve image path relative to markdown file directory"""
        # Check if it's a URL (http:// or https://)
        if image_path.startswith('http://') or image_path.startswith('https://'):
            return image_path
        
        # If already absolute and exists, use it
        if os.path.isabs(image_path) and os.path.exists(image_path):
            return image_path
        
        # Try relative to markdown directory
        resolved = os.path.join(self.md_dir, image_path)
        if os.path.exists(resolved):
            return resolved
        
        # Try as-is
        if os.path.exists(image_path):
            return image_path
        
        # Return original if nothing found (will be handled by add_image)
        return image_path
    
    def generate(self):
        """Generate PowerPoint from Markdown"""
        # Read Markdown
        with open(self.md_file, 'r', encoding='utf-8') as f:
            md_text = f.read()
        
        # Parse Markdown
        self.parser = MarkdownParser(md_text)
        slides_data = self.parser.parse()
        
        # Generate slides
        for i, slide_data in enumerate(slides_data):
            slide_type = self.parser.detect_slide_type(slide_data)
            # print(f"Slide {i+1}: Type={slide_type}, Title={slide_data.get('title')}, Images={len(slide_data.get('images', []))}, Content items={len(slide_data.get('content', []))}")
            if slide_data.get('images'):
                pass # print(f"  Image paths: {slide_data.get('images')}")
            self._create_slide(slide_data, slide_type)

        # Fix common XML issues before saving
        # This prevents PowerPoint from showing a "repair" dialog
        _fix_presentation_xml(self.prs)

        # Save presentation
        self.prs.save(self.output_file)
        # print(f'Professional PPTX presentation saved as {self.output_file}')
    
    def _create_slide(self, slide_data, slide_type):
        """Create a slide based on type"""
        slide = self.layout_engine.create_blank_slide()
        
        if slide_type == 'title_slide':
            self._create_title_slide(slide, slide_data)
        elif slide_type == 'title_slide_subtitle':
            self._create_title_slide_subtitle(slide, slide_data)
        elif slide_type == 'image_focused':
            self._create_image_focused_slide(slide, slide_data)
        elif slide_type == 'image_content':
            self._create_image_content_slide(slide, slide_data)
        elif slide_type == 'table':
            self._create_table_slide(slide, slide_data)
        elif slide_type == 'code':
            self._create_code_slide(slide, slide_data)
        elif slide_type == 'bullets':
            self._create_bullet_slide(slide, slide_data)
        elif slide_type == 'title_subtitle':
            self._create_title_subtitle_slide(slide, slide_data)
        else:
            self._create_title_content_slide(slide, slide_data)
        
        # Add footer if configured
        footer_config = LayoutEngine.get_footer_config()
        if footer_config:
            text = None
            show_on_title = True
            
            if isinstance(footer_config, dict):
                text = footer_config.get('text')
                show_on_title = footer_config.get('show_on_title_slide', True)
            else:
                text = str(footer_config)
                
            if text:
                is_title_slide = slide_type in ['title_slide', 'title_slide_subtitle']
                if not is_title_slide or show_on_title:
                    self.layout_engine.add_footer(slide, text)
    
    def _create_title_slide(self, slide, slide_data):
        """Create title slide (H1 only) with one image based on orientation"""
        if not slide_data['images']:
            return
        
        image_path = self._resolve_image_path(slide_data['images'][0])
        orientation = self.layout_engine.get_image_orientation(image_path)
        
        if orientation == 'portrait':
            # Portrait: Image scaled to full slide height on LEFT, Title on RIGHT
            img_height = self.layout_engine.SLIDE_HEIGHT
            # Calculate width maintaining aspect ratio
            with Image.open(image_path) as img:
                img_width_px, img_height_px = img.size
                aspect_ratio = img_width_px / img_height_px
            img_width = img_height * aspect_ratio
            
            # Add image on left
            self.layout_engine.add_image(slide, image_path, Inches(0), Inches(0), 
                                        width=img_width, height=img_height, maintain_aspect=False)
            
            # Add title on right
            title_left = img_width + Inches(0.5)
            title_width = self.layout_engine.SLIDE_WIDTH - title_left - self.layout_engine.MARGIN_RIGHT
            title_top = self.layout_engine.center_vertically(Inches(2))
            self.layout_engine.add_text_box(
                slide, slide_data['title'],
                title_left, title_top,
                title_width, Inches(2),
                font_size=Typography._get_size('title'),
                alignment=PP_ALIGN.CENTER
            )
            
        elif orientation == 'landscape':
            # Landscape: Fill slide with image, heading on top (z-index) centered
            # Add image first (lower z-index)
            self.layout_engine.add_image(slide, image_path, Inches(0), Inches(0),
                                        width=self.layout_engine.SLIDE_WIDTH,
                                        height=self.layout_engine.SLIDE_HEIGHT,
                                        maintain_aspect=False)
            
            # Add title on top (higher z-index, added after image)
            title_width = self.layout_engine.SLIDE_WIDTH - self.layout_engine.MARGIN_LEFT - self.layout_engine.MARGIN_RIGHT
            # Center vertically
            title_top = self.layout_engine.center_vertically(Inches(2))
            text_box = slide.shapes.add_textbox(
                self.layout_engine.MARGIN_LEFT, title_top,
                title_width, Inches(2)
            )
            text_frame = text_box.text_frame
            text_frame.text = slide_data['title']
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            Typography.apply_title_style(paragraph)
            
        else:  # square
            # Square: Title in header, image centered in body
            # Title at top
            self.layout_engine.add_title(
                slide, slide_data['title'],
                top=self.layout_engine.MARGIN_TOP,
                center=True
            )
            
            # Image centered in body
            img_size = min(self.layout_engine.CONTENT_WIDTH, self.layout_engine.CONTENT_HEIGHT - Inches(1))
            img_left = self.layout_engine.center_horizontally(img_size)
            img_top = self.layout_engine.MARGIN_TOP + Inches(1.5)
            self.layout_engine.add_image(
                slide, image_path, img_left, img_top,
                max_width=img_size, max_height=img_size
            )
    
    def _create_title_slide_subtitle(self, slide, slide_data):
        """Create title slide with subtitle (H1 and H2) with one image based on orientation"""
        if not slide_data['images']:
            return
        
        image_path = self._resolve_image_path(slide_data['images'][0])
        orientation = self.layout_engine.get_image_orientation(image_path)
        
        if orientation == 'portrait':
            # Portrait: Image scaled to full slide height on RIGHT, Title in header, subtitle on LEFT
            img_height = self.layout_engine.SLIDE_HEIGHT
            # Calculate width maintaining aspect ratio
            with Image.open(image_path) as img:
                img_width_px, img_height_px = img.size
                aspect_ratio = img_width_px / img_height_px
            img_width = img_height * aspect_ratio
            
            # Add image on right
            img_left = self.layout_engine.SLIDE_WIDTH - img_width
            self.layout_engine.add_image(slide, image_path, img_left, Inches(0),
                                        width=img_width, height=img_height, maintain_aspect=False)
            
            # Title in header (top left area)
            title_left = self.layout_engine.MARGIN_LEFT
            title_width = img_left - title_left - Inches(0.5)
            title_top = self.layout_engine.MARGIN_TOP
            title_box = self.layout_engine.add_text_box(
                slide, slide_data['title'],
                title_left, title_top,
                title_width, Inches(1.5),
                font_size=Typography._get_size('title'),
                alignment=PP_ALIGN.CENTER
            )
            
            # Subtitle on left (below title area)
            # Title box height is Inches(1.5), add spacing
            subtitle_top = title_top + Inches(1.5) + Inches(0.4)
            subtitle_width = title_width
            self.layout_engine.add_text_box(
                slide, slide_data['subtitle'],
                title_left, subtitle_top,
                subtitle_width, Inches(1),
                font_size=Typography._get_size('subtitle'),
                alignment=PP_ALIGN.CENTER
            )
            
        elif orientation == 'landscape':
            # Landscape: Fill slide with image, heading on top (z-index) centered with subtitle
            # Add image first (lower z-index)
            self.layout_engine.add_image(slide, image_path, Inches(0), Inches(0),
                                        width=self.layout_engine.SLIDE_WIDTH,
                                        height=self.layout_engine.SLIDE_HEIGHT,
                                        maintain_aspect=False)
            
            # Add title and subtitle on top (higher z-index, added after image)
            content_width = self.layout_engine.SLIDE_WIDTH - self.layout_engine.MARGIN_LEFT - self.layout_engine.MARGIN_RIGHT
            # Center content block vertically
            total_height = Inches(1.5) + Inches(1) + Inches(0.1) # Title + Subtitle + spacing
            content_top = self.layout_engine.center_vertically(total_height)
            
            # Title
            title_box = slide.shapes.add_textbox(
                self.layout_engine.MARGIN_LEFT, content_top,
                content_width, Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            title_para = title_frame.paragraphs[0]
            title_para.alignment = PP_ALIGN.CENTER
            Typography.apply_title_style(title_para)
            
            # Subtitle
            subtitle_top = content_top + Inches(1.6)
            subtitle_box = slide.shapes.add_textbox(
                self.layout_engine.MARGIN_LEFT, subtitle_top,
                content_width, Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = slide_data['subtitle']
            subtitle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            Typography.apply_body_style(subtitle_para, Typography._get_size('subtitle'))
            
        else:  # square
            # Square: Title in header, Subtitle centered in top of body, image centered in body below subtitle
            # Title at top
            title_box = self.layout_engine.add_title(
                slide, slide_data['title'],
                top=self.layout_engine.MARGIN_TOP,
                center=True
            )
            
            # Subtitle centered in top of body
            # Title box height is now Inches(1.5), add spacing
            subtitle_top = self.layout_engine.MARGIN_TOP + Inches(1.5) + Inches(0.4)
            self.layout_engine.add_text_box(
                slide, slide_data['subtitle'],
                self.layout_engine.CONTENT_LEFT, subtitle_top,
                self.layout_engine.CONTENT_WIDTH, Inches(0.75),
                font_size=Typography._get_size('subtitle'),
                alignment=PP_ALIGN.CENTER
            )
            
            # Image centered in body below subtitle
            img_size = min(self.layout_engine.CONTENT_WIDTH, Inches(4))
            img_left = self.layout_engine.center_horizontally(img_size)
            img_top = subtitle_top + Inches(1)
            self.layout_engine.add_image(
                slide, image_path, img_left, img_top,
                max_width=img_size, max_height=img_size
            )
    
    def _create_image_focused_slide(self, slide, slide_data):
        """Create slide with large centered image(s)"""
        has_title = bool(slide_data.get('title'))
        has_subtitle = bool(slide_data.get('subtitle'))
        
        # Title (optional)
        title_bottom = self.layout_engine.MARGIN_TOP
        if has_title:
            self.layout_engine.add_title(
                slide, slide_data['title'], 
                top=self.layout_engine.MARGIN_TOP, 
                center=True
            )
            title_bottom = self.layout_engine.MARGIN_TOP + Inches(1.0)
        
        # Subtitle (optional)
        if has_subtitle:
            subtitle_top = title_bottom + Inches(0.1)
            self.layout_engine.add_text_box(
                slide, slide_data['subtitle'],
                self.layout_engine.CONTENT_LEFT, subtitle_top,
                self.layout_engine.CONTENT_WIDTH, Inches(0.5),
                font_size=Typography._get_size('subtitle'),
                alignment=PP_ALIGN.CENTER
            )
        
        # Large centered image(s)
        images = slide_data.get('images', [])
        if images:
            num_images = len(images)
            if num_images == 1:
                # Single image
                image_path = self._resolve_image_path(images[0])
                
                if has_title:
                    # With title: use constrained size and position below title
                    img_width = self.layout_engine._get_image_max('max_width')
                    img_height = self.layout_engine._get_image_max('max_height')
                    left = self.layout_engine.center_horizontally(img_width)
                    top = Inches(2.5)
                    
                    self.layout_engine.add_image(
                        slide, image_path, left, top, 
                        max_width=img_width, max_height=img_height
                    )
                else:
                    # No title: fill the entire slide
                    self.layout_engine.add_image(
                        slide, image_path, 
                        left=Inches(0), top=Inches(0),
                        width=self.layout_engine.SLIDE_WIDTH,
                        height=self.layout_engine.SLIDE_HEIGHT,
                        maintain_aspect=True
                    )
            else:
                # Multiple images - arrange side by side
                if has_title:
                    img_width = Inches(4.0)
                    img_height = Inches(3.0)
                    top = Inches(2.5)
                else:
                    # No title: use larger images to fill more of the slide
                    img_width = Inches(5.5)
                    img_height = Inches(4.0)
                    top = Inches(1.75)
                
                img_spacing = Inches(0.5)
                total_width = (img_width * num_images) + (img_spacing * (num_images - 1))
                start_left = self.layout_engine.center_horizontally(total_width)
                
                for i, img_path in enumerate(images):
                    image_path = self._resolve_image_path(img_path)
                    img_left = start_left + (img_width + img_spacing) * i
                    self.layout_engine.add_image(
                        slide, image_path, img_left, top,
                        width=img_width, height=img_height
                    )
    
    def _create_image_content_slide(self, slide, slide_data):
        """Create slide with title, subtitle, content, and image(s)"""
        images = slide_data.get('images', [])
        # print(f"Creating image_content slide - Images: {images}, Content: {len(slide_data.get('content', []))} items")
        
        # Title
        title_bottom = self.layout_engine.MARGIN_TOP
        if slide_data.get('title'):
            self.layout_engine.add_title(
                slide, slide_data['title'], 
                top=self.layout_engine.MARGIN_TOP, 
                center=False
            )
            title_bottom = self.layout_engine.MARGIN_TOP + Inches(1.0)
        
        # Subtitle
        content_top = Inches(1.5)
        if slide_data.get('subtitle'):
            subtitle_top = title_bottom + Inches(0.1)
            self.layout_engine.add_text_box(
                slide, slide_data['subtitle'],
                self.layout_engine.CONTENT_LEFT, subtitle_top,
                self.layout_engine.CONTENT_WIDTH, Inches(0.5),
                font_size=Typography._get_size('subtitle'),
                alignment=PP_ALIGN.LEFT
            )
            content_top = subtitle_top + Inches(0.6)
        
        # Two-column layout: content left, image right
        content_left = self.layout_engine.CONTENT_LEFT
        content_width = Inches(6.5)
        content_height = Inches(5.5)
        
        # Prepare a single text box for all textual content (auto-sized)
        text_box = None
        text_frame = None
        first_paragraph = True
        
        def ensure_text_frame():
            nonlocal text_box, text_frame, first_paragraph
            if text_box is None:
                text_box = slide.shapes.add_textbox(
                    content_left, content_top, content_width, content_height
                )
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                text_frame.vertical_anchor = MSO_ANCHOR.TOP
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                text_frame.clear()
                # ensure there is at least one paragraph to work with
                _ = text_frame.paragraphs[0]
                first_paragraph = True
            return text_frame
        
        def add_new_paragraph():
            frame = ensure_text_frame()
            nonlocal first_paragraph
            if first_paragraph:
                paragraph = frame.paragraphs[0]
                paragraph.text = ""
                first_paragraph = False
            else:
                paragraph = frame.add_paragraph()
                paragraph.text = ""
            paragraph.alignment = PP_ALIGN.LEFT
            return paragraph
        
        # Add content (bullets or paragraphs) - if no content, images will be on left
        has_content = False
        for content_item in slide_data.get('content', []):
            has_content = True
            if content_item['type'] == 'bullets':
                items = content_item['items'] # Pass full items with level
                last_bullet_para = None
                for item in items:
                    if isinstance(item, dict) and 'text' in item:
                        text = item['text']
                        level = item.get('level', 0)
                    else:
                        text = str(item)
                        level = 0
                    
                    p = add_new_paragraph()
                    p.level = level
                    p.font.name = Typography._get_font('body')
                    p.font.size = Typography._get_size('bullet')
                    p.font.color.rgb = Typography._get_color('primary')
                    p.space_after = Pt(Typography._get_spacing('bullet_space_after'))
                    
                    # Add bullet symbol
                    run = p.add_run()
                    run.text = "• "
                    
                    # Add formatted text
                    self.layout_engine._process_text_with_formatting(p, text)
                    last_bullet_para = p
                
                # Add extra spacing after bullet section
                if last_bullet_para:
                    last_bullet_para.space_after = Pt(
                        Typography._get_spacing('bullet_space_after') + 6
                    )
                
            elif content_item['type'] == 'paragraph':
                text = content_item['text']
                p = add_new_paragraph()
                Typography.apply_body_style(p)
                self.layout_engine._process_text_with_formatting(p, text)
            
            elif content_item['type'] == 'numbered':
                items = content_item['items']
                last_numbered_para = None
                for i, item in enumerate(items, start=1):
                    if isinstance(item, dict) and 'text' in item:
                        text = item['text']
                    else:
                        text = str(item)
                    
                    p = add_new_paragraph()
                    p.font.name = Typography._get_font('body')
                    p.font.size = Typography._get_size('bullet')
                    p.font.color.rgb = Typography._get_color('primary')
                    p.space_after = Pt(Typography._get_spacing('bullet_space_after'))
                    
                    # Add number prefix
                    run = p.add_run()
                    run.text = f"{i}. "
                    
                    # Add formatted text
                    self.layout_engine._process_text_with_formatting(p, text)
                    last_numbered_para = p
                
                # Add extra spacing after numbered list section
                if last_numbered_para:
                    last_numbered_para.space_after = Pt(
                        Typography._get_spacing('bullet_space_after') + 6
                    )
            
            else:
                # Fallback: treat as paragraph
                text = content_item.get('text', '')
                if text:
                    p = add_new_paragraph()
                    Typography.apply_body_style(p)
                    self.layout_engine._process_text_with_formatting(p, text)
        
        # Add images - handle multiple images
        if images:
            num_images = len(images)
            # print(f"Adding {num_images} image(s)")
            
            if not has_content:
                # No content, arrange images in a grid or side-by-side
                if num_images == 1:
                    # Single image, center it
                    image_path = self._resolve_image_path(images[0])
                    img_width = self.layout_engine._get_image_max('max_width')
                    img_height = self.layout_engine._get_image_max('max_height')
                    img_left = self.layout_engine.center_horizontally(img_width)
                    img_top = Inches(2.5)
                    result = self.layout_engine.add_image(
                        slide, image_path, img_left, img_top,
                        max_width=img_width, max_height=img_height
                    )
                    if result is None:
                        print(f"WARNING: Failed to add image {image_path}")
                else:
                    # Multiple images - arrange side by side or in grid
                    img_width = Inches(4.0)
                    img_height = Inches(3.0)
                    img_spacing = Inches(0.5)
                    total_width = (img_width * num_images) + (img_spacing * (num_images - 1))
                    start_left = self.layout_engine.center_horizontally(total_width)
                    img_top = Inches(2.5)
                    
                    for i, img_path in enumerate(images):
                        image_path = self._resolve_image_path(img_path)
                        img_left = start_left + (img_width + img_spacing) * i
                        # print(f"Adding image {i+1}/{num_images}: {image_path}")
                        result = self.layout_engine.add_image(
                            slide, image_path, img_left, img_top,
                            width=img_width, height=img_height
                        )
                        if result is None:
                            print(f"WARNING: Failed to add image {image_path}")
            else:
                # Content exists, put images on right
                if num_images == 1:
                    # Single image on right
                    image_path = self._resolve_image_path(images[0])
                    img_left = Inches(8)
                    img_top = Inches(1.5)
                    img_width = Inches(4.5)
                    img_height = content_height
                    # print(f"Adding image: {image_path}")
                    result = self.layout_engine.add_image(
                        slide, image_path, img_left, img_top,
                        max_width=img_width, max_height=img_height
                    )
                    if result is None:
                        print(f"WARNING: Failed to add image {image_path}")
                else:
                    # Multiple images - stack them vertically on the right
                    img_width = Inches(4.5)
                    img_height_per_image = content_height / num_images
                    img_spacing = Inches(0.2)
                    img_left = Inches(8)
                    img_top = Inches(1.5)
                    
                    for i, img_path in enumerate(images):
                        image_path = self._resolve_image_path(img_path)
                        current_top = img_top + (img_height_per_image + img_spacing) * i
                        # print(f"Adding image {i+1}/{num_images}: {image_path}")
                        result = self.layout_engine.add_image(
                            slide, image_path, img_left, current_top,
                            max_width=img_width, max_height=img_height_per_image
                        )
                        if result is None:
                            print(f"WARNING: Failed to add image {image_path}")
        else:
            print("WARNING: No images found in slide_data")
    
    def _create_table_slide(self, slide, slide_data):
        """Create slide with table"""
        # Title
        if slide_data.get('title'):
            self.layout_engine.add_title(
                slide, slide_data['title'], 
                top=self.layout_engine.MARGIN_TOP, 
                center=False
            )
        
        # Find table in content
        table_data = None
        for content_item in slide_data.get('content', []):
            if content_item['type'] == 'table':
                table_data = content_item
                break
        
        if not table_data:
            return
        
        headers = table_data.get('headers')
        rows = table_data.get('rows', [])
        
        if not rows:
            return
        
        # Determine table dimensions
        num_cols = len(rows[0]) if rows else (len(headers) if headers else 0)
        if headers:
            num_cols = max(num_cols, len(headers))
        num_rows = len(rows) + (1 if headers else 0)
        
        if num_cols == 0 or num_rows == 0:
            return
        
        # Check if table has images
        has_images = False
        if headers:
            for h in headers:
                if isinstance(h, dict) and h.get('images'):
                    has_images = True
                    break
        if not has_images:
            for r in rows:
                for c in r:
                    if isinstance(c, dict) and c.get('images'):
                        has_images = True
                        break
                if has_images:
                    break
        
        # Calculate table position and size
        table_top = Inches(1.5) if slide_data.get('title') else Inches(1.0)
        table_left = self.layout_engine.CONTENT_LEFT
        table_width = self.layout_engine.CONTENT_WIDTH
        
        # Determine appropriate height
        max_height = min(Inches(5.5), self.layout_engine.SLIDE_HEIGHT - table_top - self.layout_engine.MARGIN_BOTTOM)
        
        if has_images:
            # For tables with images, use more space to accommodate them
            table_height = max_height
        else:
            # For text-only tables, scale height by row count to avoid excessively tall cells
            # Base height per row approx 0.6 inches
            calculated_height = Inches(0.6 * num_rows)
            table_height = min(calculated_height, max_height)
        
        # Add table shape
        table_shape = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height)
        table = table_shape.table
        
        # Set column widths
        col_width = table_width / num_cols
        for col_idx in range(num_cols):
            table.columns[col_idx].width = int(col_width)
        
        # Helper function to add image to cell
        def add_image_to_cell(cell, image_path, cell_left, cell_top, cell_width, cell_height):
            """Add an image to a table cell by positioning it within the cell bounds"""
            try:
                # Resolve image path (handles URLs and local files)
                resolved_path = self._resolve_image_path(image_path)
                
                # Calculate image size to fit in cell with some padding
                padding = Inches(0.1)
                max_img_width = cell_width - (padding * 2)
                max_img_height = cell_height - (padding * 2)
                
                # Get image dimensions
                try:
                    is_url = resolved_path.startswith('http://') or resolved_path.startswith('https://')
                    temp_file = None
                    actual_path = resolved_path
                    
                    if is_url:
                        # Download URL to temp file with proper headers
                        url_request = _create_url_request(resolved_path)
                        with urlopen(url_request) as req:
                            img_data = req.read()
                            if img_data:
                                content_type = req.headers.get('Content-Type', 'image/jpeg')
                                suffix = '.png' if 'png' in content_type.lower() else '.jpg'
                                
                                # Convert unsupported formats (like WEBP) on the fly
                                try:
                                    from io import BytesIO
                                    with Image.open(BytesIO(img_data)) as img:
                                        if img.format and img.format.upper() not in PPTX_SUPPORTED_FORMATS:
                                            output = BytesIO()
                                            if img.mode in ('RGBA', 'LA', 'P'):
                                                if img.mode == 'P' and 'transparency' in img.info:
                                                    img = img.convert('RGBA')
                                                elif img.mode == 'P':
                                                    img = img.convert('RGB')
                                            elif img.mode not in ('RGB', 'RGBA'):
                                                img = img.convert('RGB')
                                            img.save(output, 'PNG')
                                            img_data = output.getvalue()
                                            suffix = '.png'
                                except Exception as conv_err:
                                    pass  # Continue with original format
                                
                                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
                                temp_file.write(img_data)
                                temp_file.close()
                                actual_path = temp_file.name
                    
                    with Image.open(actual_path) as img:
                        img_width_px, img_height_px = img.size
                        aspect_ratio = img_width_px / img_height_px
                    
                    # Calculate size maintaining aspect ratio
                    if max_img_width / aspect_ratio <= max_img_height:
                        img_width = int(max_img_width)
                        img_height = int(max_img_width / aspect_ratio)
                    else:
                        img_height = int(max_img_height)
                        img_width = int(max_img_height * aspect_ratio)
                    
                    # Center image in cell
                    img_left = int(cell_left + (cell_width - img_width) / 2)
                    img_top = int(cell_top + (cell_height - img_height) / 2)
                    
                    # For local files, convert if needed; URLs already converted on download
                    if not is_url:
                        converted_path, needs_conversion_cleanup = _convert_to_supported_format(actual_path)
                        if needs_conversion_cleanup:
                            actual_path = converted_path
                    else:
                        needs_conversion_cleanup = False
                    
                    # Add picture to slide (positioned over the cell)
                    picture = slide.shapes.add_picture(actual_path, img_left, img_top, 
                                                     width=img_width, height=img_height)
                    
                    # Clean up converted temp file if we created one (for local files)
                    if needs_conversion_cleanup:
                        try:
                            os.unlink(actual_path)
                        except:
                            pass
                    
                    # Clean up URL temp file if created
                    if temp_file:
                        try:
                            os.unlink(temp_file.name)
                        except:
                            pass
                    
                    return picture
                except Exception as e:
                    print(f"Error processing image {image_path} for table cell: {e}")
                    if temp_file:
                        try:
                            os.unlink(temp_file.name)
                        except:
                            pass
                    return None
            except Exception as e:
                print(f"Error adding image to table cell: {e}")
                return None
        
        # Add headers if present
        row_idx = 0
        if headers:
            for col_idx, header_cell in enumerate(headers):
                if col_idx < num_cols:
                    cell = table.cell(row_idx, col_idx)
                    header_text = header_cell.get('text', '') if isinstance(header_cell, dict) else str(header_cell)
                    header_images = header_cell.get('images', []) if isinstance(header_cell, dict) else []
                    
                    # Only add text if no images (or add both if needed)
                    if not header_images:
                        cell.text = header_text
                    
                    # Style header
                    if cell.text_frame.paragraphs:
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.bold = True
                        paragraph.font.size = Typography._get_size('body')
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
                        paragraph.alignment = PP_ALIGN.CENTER
                    
                    # Header background
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = Typography._get_color('primary')
                    
                    # Add images to header cell if present
                    if header_images:
                        # Get cell position - calculate from table position and column/row dimensions
                        # All values in EMUs (integers)
                        cell_left = table_left + sum(table.columns[i].width for i in range(col_idx))
                        # Row heights are not directly accessible easily, so we estimate evenly distributed
                        cell_top = table_top + int((row_idx * table_height) / num_rows)
                        cell_width = table.columns[col_idx].width
                        cell_height = int(table_height / num_rows)
                        
                        for img_path in header_images:
                            add_image_to_cell(cell, img_path, cell_left, cell_top, cell_width, cell_height)
            row_idx += 1
        
        # Add data rows
        for row_data in rows:
            if row_idx >= num_rows:
                break
            for col_idx, cell_data in enumerate(row_data):
                if col_idx < num_cols:
                    cell = table.cell(row_idx, col_idx)
                    cell_text = cell_data.get('text', '') if isinstance(cell_data, dict) else str(cell_data)
                    cell_images = cell_data.get('images', []) if isinstance(cell_data, dict) else []
                    
                    # Only add text if no images (or add both if needed)
                    if not cell_images:
                        cell.text = cell_text
                    
                    # Style data cell
                    if cell.text_frame.paragraphs:
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Typography._get_size('body')
                        paragraph.font.color.rgb = Typography._get_color('primary')
                        paragraph.alignment = PP_ALIGN.LEFT
                    
                    # Alternating row background (only if no images, as images will cover it)
                    if row_idx % 2 == 0 and not cell_images:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = Typography._get_color('light_bg')
                    
                    # Add images to cell if present
                    if cell_images:
                        # Get cell position - calculate from table position and column/row dimensions
                        # All values in EMUs (integers)
                        cell_left = table_left + sum(table.columns[i].width for i in range(col_idx))
                        # Row heights are not directly accessible easily, so we estimate evenly distributed
                        cell_top = table_top + int((row_idx * table_height) / num_rows)
                        cell_width = table.columns[col_idx].width
                        cell_height = int(table_height / num_rows)
                        
                        for img_path in cell_images:
                            add_image_to_cell(cell, img_path, cell_left, cell_top, cell_width, cell_height)
            row_idx += 1
    
    def _create_code_slide(self, slide, slide_data):
        """Create slide with code block"""
        # Title
        if slide_data['title']:
            self.layout_engine.add_title(
                slide, slide_data['title'], 
                top=self.layout_engine.MARGIN_TOP, 
                center=False
            )
        
        # Find code block
        code_text = None
        for content_item in slide_data['content']:
            if content_item['type'] == 'code':
                code_text = content_item['text']
                break
        
        if code_text:
            # Create code block with background
            code_left = self.layout_engine.CONTENT_LEFT
            code_top = Inches(1.5)
            code_width = self.layout_engine.CONTENT_WIDTH
            code_height = Inches(5)
            
            # Add background shape
            code_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, code_left, code_top, 
                code_width, code_height
            )
            code_box.fill.solid()
            code_box.fill.fore_color.rgb = Typography._get_color('code_bg')
            code_box.line.color.rgb = Typography._get_color('code_border')
            code_box.line.width = Pt(1)
            
            # Add text
            text_frame = code_box.text_frame
            text_frame.text = code_text
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.25)
            text_frame.margin_right = Inches(0.25)
            text_frame.margin_top = Inches(0.25)
            text_frame.margin_bottom = Inches(0.25)
            
            paragraph = text_frame.paragraphs[0]
            Typography.apply_code_style(paragraph)
    
    def _create_bullet_slide(self, slide, slide_data):
        """Create slide with title, subtitle, and bullet list(s), handling mixed content"""
        # Title
        title_bottom = self.layout_engine.MARGIN_TOP
        if slide_data.get('title'):
            self.layout_engine.add_title(
                slide, slide_data['title'], 
                top=self.layout_engine.MARGIN_TOP, 
                center=False
            )
            title_bottom = self.layout_engine.MARGIN_TOP + Inches(1.0)
        
        # Subtitle
        content_top = Inches(1.5)
        if slide_data.get('subtitle'):
            subtitle_top = title_bottom + Inches(0.1)
            self.layout_engine.add_text_box(
                slide, slide_data['subtitle'],
                self.layout_engine.CONTENT_LEFT, subtitle_top,
                self.layout_engine.CONTENT_WIDTH, Inches(0.5),
                font_size=Typography._get_size('subtitle'),
                alignment=PP_ALIGN.LEFT
            )
            content_top = subtitle_top + Inches(0.6)
        
        # Create a single text box for all content (handles mixed paragraphs and bullet lists)
        content_left = self.layout_engine.CONTENT_LEFT
        content_width = self.layout_engine.CONTENT_WIDTH
        content_height = Inches(5.5)
        
        text_box = slide.shapes.add_textbox(
            content_left, content_top, content_width, content_height
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.clear()
        
        # Ensure there's at least one paragraph to work with
        _ = text_frame.paragraphs[0]
        first_paragraph = True
        
        def add_new_paragraph():
            nonlocal first_paragraph
            if first_paragraph:
                paragraph = text_frame.paragraphs[0]
                paragraph.text = ""
                first_paragraph = False
            else:
                paragraph = text_frame.add_paragraph()
                paragraph.text = ""
            paragraph.alignment = PP_ALIGN.LEFT
            return paragraph
        
        # Process all content items in order (paragraphs, bullet lists, and numbered lists)
        for content_item in slide_data['content']:
            if content_item['type'] == 'bullets':
                items = content_item['items']
                last_bullet_para = None
                for item in items:
                    if isinstance(item, dict) and 'text' in item:
                        text = item['text']
                        level = item.get('level', 0)
                    else:
                        text = str(item)
                        level = 0
                    
                    p = add_new_paragraph()
                    p.level = level
                    p.font.name = Typography._get_font('body')
                    p.font.size = Typography._get_size('bullet')
                    p.font.color.rgb = Typography._get_color('primary')
                    p.space_after = Pt(Typography._get_spacing('bullet_space_after'))
                    
                    # Add bullet symbol
                    run = p.add_run()
                    run.text = "• "
                    
                    # Add formatted text
                    self.layout_engine._process_text_with_formatting(p, text)
                    last_bullet_para = p
                
                # Add extra spacing after bullet section
                if last_bullet_para:
                    last_bullet_para.space_after = Pt(
                        Typography._get_spacing('bullet_space_after') + 6
                    )
            
            elif content_item['type'] == 'numbered':
                items = content_item['items']
                last_numbered_para = None
                for i, item in enumerate(items, start=1):
                    if isinstance(item, dict) and 'text' in item:
                        text = item['text']
                    else:
                        text = str(item)
                    
                    p = add_new_paragraph()
                    p.font.name = Typography._get_font('body')
                    p.font.size = Typography._get_size('bullet')
                    p.font.color.rgb = Typography._get_color('primary')
                    p.space_after = Pt(Typography._get_spacing('bullet_space_after'))
                    
                    # Add number prefix
                    run = p.add_run()
                    run.text = f"{i}. "
                    
                    # Add formatted text
                    self.layout_engine._process_text_with_formatting(p, text)
                    last_numbered_para = p
                
                # Add extra spacing after numbered list section
                if last_numbered_para:
                    last_numbered_para.space_after = Pt(
                        Typography._get_spacing('bullet_space_after') + 6
                    )
            
            elif content_item['type'] == 'paragraph':
                text = content_item['text']
                p = add_new_paragraph()
                Typography.apply_body_style(p)
                self.layout_engine._process_text_with_formatting(p, text)
                # Add spacing after paragraph
                p.space_after = Pt(12)
            
            elif content_item['type'] == 'heading3':
                text = content_item['text']
                p = add_new_paragraph()
                p.font.name = Typography._get_font('body')
                p.font.size = Typography._get_size('body') + Pt(2)  # Slightly larger than body
                p.font.color.rgb = Typography._get_color('primary')
                p.font.bold = True
                p.text = text
                p.space_after = Pt(6)
            
            else:
                # Fallback: treat as paragraph
                text = content_item.get('text', '')
                if text:
                    p = add_new_paragraph()
                    Typography.apply_body_style(p)
                    self.layout_engine._process_text_with_formatting(p, text)
                    p.space_after = Pt(12)
    
    def _create_title_subtitle_slide(self, slide, slide_data):
        """Create slide with title and subtitle"""
        # Title
        title_box = None
        if slide_data['title']:
            # Calculate total height needed: title (1.5") + spacing (0.4") + subtitle (0.75")
            total_height = Inches(1.5) + Inches(0.4) + Inches(0.75)
            title_top = self.layout_engine.center_vertically(total_height)
            title_box = self.layout_engine.add_title(
                slide, slide_data['title'], 
                top=title_top, 
                center=True
            )
        
        # Subtitle
        if slide_data['subtitle']:
            # Calculate subtitle position based on title box bottom + spacing
            if title_box:
                # Title box height is now Inches(1.5), so bottom is title_top + 1.5
                # Add spacing after title
                subtitle_top = title_top + Inches(1.5) + Inches(0.4)
            else:
                # No title, center subtitle
                subtitle_top = self.layout_engine.center_vertically(Inches(1))
            
            self.layout_engine.add_text_box(
                slide, slide_data['subtitle'],
                self.layout_engine.CONTENT_LEFT, subtitle_top,
                self.layout_engine.CONTENT_WIDTH, Inches(0.75),
                font_size=Typography._get_size('subtitle'),
                alignment=PP_ALIGN.CENTER
            )
    
    def _create_title_content_slide(self, slide, slide_data):
        """Create standard title + subtitle + content slide with proper paragraph handling"""
        # Title
        title_bottom = self.layout_engine.MARGIN_TOP
        if slide_data.get('title'):
            self.layout_engine.add_title(
                slide, slide_data['title'], 
                top=self.layout_engine.MARGIN_TOP, 
                center=False
            )
            title_bottom = self.layout_engine.MARGIN_TOP + Inches(1.0)
        
        # Subtitle
        content_top = Inches(1.5)
        if slide_data.get('subtitle'):
            subtitle_top = title_bottom + Inches(0.1)
            self.layout_engine.add_text_box(
                slide, slide_data['subtitle'],
                self.layout_engine.CONTENT_LEFT, subtitle_top,
                self.layout_engine.CONTENT_WIDTH, Inches(0.5),
                font_size=Typography._get_size('subtitle'),
                alignment=PP_ALIGN.LEFT
            )
            content_top = subtitle_top + Inches(0.6)
        
        # Create a single text box for all content (handles multiple paragraphs properly)
        content_left = self.layout_engine.CONTENT_LEFT
        content_width = self.layout_engine.CONTENT_WIDTH
        content_height = Inches(5.5)
        
        text_box = slide.shapes.add_textbox(
            content_left, content_top, content_width, content_height
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.clear()
        
        # Ensure there's at least one paragraph to work with
        _ = text_frame.paragraphs[0]
        first_paragraph = True
        
        def add_new_paragraph():
            nonlocal first_paragraph
            if first_paragraph:
                paragraph = text_frame.paragraphs[0]
                paragraph.text = ""
                first_paragraph = False
            else:
                paragraph = text_frame.add_paragraph()
                paragraph.text = ""
            paragraph.alignment = PP_ALIGN.LEFT
            return paragraph
        
        # Process all content items in order
        for content_item in slide_data['content']:
            if content_item['type'] == 'paragraph':
                text = content_item['text']
                p = add_new_paragraph()
                Typography.apply_body_style(p)
                self.layout_engine._process_text_with_formatting(p, text)
                # Add spacing after paragraph
                p.space_after = Pt(12)
            
            elif content_item['type'] == 'heading3':
                text = content_item['text']
                p = add_new_paragraph()
                p.font.name = Typography._get_font('body')
                p.font.size = Typography._get_size('body') + Pt(2)  # Slightly larger than body
                p.font.color.rgb = Typography._get_color('primary')
                p.font.bold = True
                p.text = text
                p.space_after = Pt(6)
            
            elif content_item['type'] == 'bullets':
                items = content_item['items']
                last_bullet_para = None
                for item in items:
                    if isinstance(item, dict) and 'text' in item:
                        text = item['text']
                        level = item.get('level', 0)
                    else:
                        text = str(item)
                        level = 0
                    
                    p = add_new_paragraph()
                    p.level = level
                    p.font.name = Typography._get_font('body')
                    p.font.size = Typography._get_size('bullet')
                    p.font.color.rgb = Typography._get_color('primary')
                    p.space_after = Pt(Typography._get_spacing('bullet_space_after'))
                    
                    # Add bullet symbol
                    run = p.add_run()
                    run.text = "• "
                    
                    # Add formatted text
                    self.layout_engine._process_text_with_formatting(p, text)
                    last_bullet_para = p
                
                # Add extra spacing after bullet section
                if last_bullet_para:
                    last_bullet_para.space_after = Pt(
                        Typography._get_spacing('bullet_space_after') + 6
                    )
            
            elif content_item['type'] == 'numbered':
                items = content_item['items']
                last_numbered_para = None
                for i, item in enumerate(items, start=1):
                    if isinstance(item, dict) and 'text' in item:
                        text = item['text']
                    else:
                        text = str(item)
                    
                    p = add_new_paragraph()
                    p.font.name = Typography._get_font('body')
                    p.font.size = Typography._get_size('bullet')
                    p.font.color.rgb = Typography._get_color('primary')
                    p.space_after = Pt(Typography._get_spacing('bullet_space_after'))
                    
                    # Add number prefix
                    run = p.add_run()
                    run.text = f"{i}. "
                    
                    # Add formatted text
                    self.layout_engine._process_text_with_formatting(p, text)
                    last_numbered_para = p
                
                # Add extra spacing after numbered list section
                if last_numbered_para:
                    last_numbered_para.space_after = Pt(
                        Typography._get_spacing('bullet_space_after') + 6
                    )
            
            else:
                # Fallback: treat as paragraph
                text = content_item.get('text', '')
                if text:
                    p = add_new_paragraph()
                    Typography.apply_body_style(p)
                    self.layout_engine._process_text_with_formatting(p, text)
                    p.space_after = Pt(12)


def main():
    """Main entry point"""
    md_file = 'slide_show_guide.md'
    output_file = 'slide_show_guide_professional.pptx'
    
    # Look for theme file in same directory as script
    script_dir = os.path.dirname(os.path.abspath(__file__))
    theme_file = os.path.join(script_dir, 'pptx_theme.json')
    
    generator = PowerPointGenerator(md_file, output_file, theme_file=theme_file)
    generator.generate()


# if __name__ == '__main__':
#    main()
